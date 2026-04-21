"""Build the Accutite 'How to Use ChatGPT in the Fastener Industry' deck.

Outputs Accutite_ChatGPT_Fastener_Deck.pptx in this directory.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------- theme ----------
NAVY_900 = RGBColor(0x04, 0x06, 0x0F)
NAVY_800 = RGBColor(0x07, 0x0B, 0x1C)
NAVY_700 = RGBColor(0x0C, 0x12, 0x30)
NAVY_600 = RGBColor(0x12, 0x1A, 0x44)
VIOLET   = RGBColor(0x7C, 0x4D, 0xFF)
VIOLET_L = RGBColor(0x9A, 0x7B, 0xFF)
CYAN     = RGBColor(0x38, 0xE0, 0xFF)
CYAN_L   = RGBColor(0x7E, 0xF0, 0xFF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INK_200  = RGBColor(0xC9, 0xD4, 0xF5)
INK_300  = RGBColor(0x8E, 0xA0, 0xD6)
INK_400  = RGBColor(0x5D, 0x6E, 0xA8)

FONT = "Calibri"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- helpers ----------
def set_background(slide, fill_rgb, second_rgb=None):
    bg = slide.background
    fill = bg.fill
    if second_rgb is None:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.gradient()
        stops = fill.gradient_stops
        stops[0].color.rgb = fill_rgb
        stops[1].color.rgb = second_rgb
        fill.gradient_angle = 90.0


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shadow=False, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    return shp


def add_text(slide, x, y, w, h, text, *, font=FONT, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, tracking=None, italic=False,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if tracking is not None:
            _set_spacing(run, tracking)
    return tb


def _set_spacing(run, hundredths_of_point):
    """Set character spacing in hundredths of a point via raw XML."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_of_point)))


def add_kicker(slide, x, y, text):
    """Pill-shaped cyan kicker label."""
    w = Inches(2.2)
    h = Inches(0.36)
    pill = add_rect(slide, x, y, w, h, fill=NAVY_800, line=CYAN, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pill.adjustments[0] = 0.5
    tb = add_text(slide, x, y, w, h, text, size=9, color=CYAN_L, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=300)
    return tb


def add_header(slide, kicker_text, title_text):
    add_kicker(slide, Inches(0.7), Inches(0.55), kicker_text)
    add_text(slide, Inches(0.7), Inches(1.0), Inches(12), Inches(1.1), title_text,
             size=36, bold=True, color=WHITE, line_spacing=1.05)
    # accent underline
    add_rect(slide, Inches(0.7), Inches(2.0), Inches(0.9), Emu(38100), fill=CYAN)


def add_footer(slide):
    y = Inches(7.05)
    add_rect(slide, Inches(0.7), y - Emu(38100), Inches(11.933), Emu(9525), fill=VIOLET_L)
    add_text(slide, Inches(0.7), y, Inches(6), Inches(0.3),
             "ACCUTITE  •  HOW TO USE CHATGPT IN THE FASTENER INDUSTRY",
             size=8, color=INK_400, bold=True, tracking=200)
    add_text(slide, Inches(6.7), y, Inches(6.0), Inches(0.3),
             "ACCUTITE FASTENERS, INC.", size=8, color=INK_400, bold=True,
             tracking=200, align=PP_ALIGN.RIGHT)


def add_page_num(slide, idx, total=16):
    add_text(slide, Inches(12.0), Inches(0.35), Inches(1.1), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=9, color=INK_400, bold=True,
             align=PP_ALIGN.RIGHT, tracking=300)


def add_cosmic_accents(slide):
    """Corner glow shapes to suggest the 'galactic' feel."""
    # top-right violet glow
    glow1 = add_rect(slide, Inches(9.5), Inches(-1.0), Inches(5.0), Inches(4.0),
                     fill=VIOLET, shape=MSO_SHAPE.OVAL)
    glow1.fill.transparency = 0.85
    glow1.line.fill.background()
    # bottom-left cyan glow
    glow2 = add_rect(slide, Inches(-1.5), Inches(5.5), Inches(5.0), Inches(3.5),
                     fill=CYAN, shape=MSO_SHAPE.OVAL)
    glow2.fill.transparency = 0.9
    glow2.line.fill.background()


def add_logo_lockup(slide, x, y, *, large=False):
    w = Inches(2.4 if large else 1.8)
    h = Inches(0.7 if large else 0.55)
    add_rect(slide, x, y, w, h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(0.75),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x, y + Emu(25400), w, h / 2, "ACCUTITE",
             size=14 if large else 11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=400)
    add_text(slide, x, y + h / 2, w, h / 2, "FASTENERS INC.",
             size=8 if large else 7, bold=True, color=INK_300,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, tracking=500)


def add_portrait_placeholder(slide, x, y, w, h, label):
    add_rect(slide, x, y, w, h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(1.25),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # inner dashed frame
    inner = add_rect(slide, x + Inches(0.2), y + Inches(0.2),
                     w - Inches(0.4), h - Inches(0.4),
                     fill=None, line=VIOLET_L, line_w=Pt(0.5),
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _dash_line(inner)
    add_text(slide, x, y + h / 2 - Inches(0.2), w, Inches(0.4),
             label, size=10, bold=True, color=INK_300, tracking=300,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _dash_line(shape):
    ln = shape.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")


def bullet_row(slide, x, y, w, h, text):
    add_rect(slide, x, y, w, h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(0.5),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # diamond marker
    d = Inches(0.22)
    add_rect(slide, x + Inches(0.3), y + h / 2 - d / 2, d, d,
             fill=VIOLET, shape=MSO_SHAPE.DIAMOND)
    add_text(slide, x + Inches(0.8), y, w - Inches(1.0), h, text,
             size=17, color=INK_200, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)


# ---------- slide builders ----------
def base_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY_800, NAVY_900)
    add_cosmic_accents(slide)
    return slide


def slide_title(prs):
    slide = base_slide(prs)

    add_text(slide, Inches(0.7), Inches(0.6), Inches(8), Inches(0.3),
             "ACCUTITE FASTENERS, INC.", size=11, bold=True, color=CYAN_L, tracking=500)

    add_text(slide, Inches(0.7), Inches(1.3), Inches(8.3), Inches(2.2),
             "How to Use ChatGPT in the", size=52, bold=True, color=WHITE, line_spacing=1.0)
    add_text(slide, Inches(0.7), Inches(2.3), Inches(8.3), Inches(1.2),
             "Fastener Industry", size=52, bold=True, color=VIOLET_L, line_spacing=1.0)

    # accent bar
    add_rect(slide, Inches(0.7), Inches(3.55), Inches(1.2), Emu(50800), fill=CYAN)

    add_text(slide, Inches(0.7), Inches(3.8), Inches(6), Inches(0.45),
             "Levee Hedrick", size=22, bold=True, color=WHITE)
    add_text(slide, Inches(0.7), Inches(4.25), Inches(6), Inches(0.4),
             "MATERIAL CONTROL MANAGER", size=11, bold=True, color=INK_300, tracking=400)

    # portrait placeholder (hero)
    add_portrait_placeholder(slide, Inches(9.0), Inches(1.3), Inches(3.6), Inches(5.2),
                             "HERO PORTRAIT")

    add_logo_lockup(slide, Inches(10.7), Inches(6.6), large=True)


def content_slide(prs, idx, kicker, title, bullets):
    slide = base_slide(prs)
    add_header(slide, kicker, title)

    top = Inches(2.35)
    row_h = Inches(0.72)
    gap = Inches(0.18)
    if len(bullets) >= 5:
        row_h = Inches(0.6)
        gap = Inches(0.12)

    for i, text in enumerate(bullets):
        y = top + (row_h + gap) * i
        bullet_row(slide, Inches(0.7), y, Inches(11.933), row_h, text)

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def chips_slide(prs, idx):
    slide = base_slide(prs)
    add_header(slide, "03 • SURFACE AREA", "Where ChatGPT Fits in Our Daily Workflow")

    chips = ["Emails", "Quotes", "PDFs", "Drawings", "Blueprints",
             "Quality Documents", "Purchasing Communication"]
    widths = {t: Inches(1.2 + 0.18 * len(t)) for t in chips}

    x, y = Inches(0.7), Inches(3.0)
    row_gap = Inches(0.3)
    row_h = Inches(0.8)
    max_x = Inches(12.6)
    for t in chips:
        w = widths[t]
        if x + w > max_x:
            x = Inches(0.7)
            y += row_h + row_gap
        chip = add_rect(slide, x, y, w, row_h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(1),
                        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        chip.adjustments[0] = 0.5
        add_text(slide, x, y, w, row_h, t, size=18, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w + Inches(0.25)

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def cards_slide(prs, idx, kicker, title, items):
    slide = base_slide(prs)
    add_header(slide, kicker, title)

    cols = 2
    card_w = Inches(5.85)
    card_h = Inches(1.6)
    gap_x = Inches(0.2)
    gap_y = Inches(0.25)
    top = Inches(2.6)

    for i, text in enumerate(items):
        r, c = divmod(i, cols)
        x = Inches(0.7) + (card_w + gap_x) * c
        y = top + (card_h + gap_y) * r
        add_rect(slide, x, y, card_w, card_h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(1),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(1.0), Inches(0.4),
                 f"{i+1:02d}", size=11, bold=True, color=CYAN_L, tracking=400)
        # accent bar inside card
        add_rect(slide, x + Inches(0.3), y + Inches(0.65), Inches(0.5), Emu(25400), fill=VIOLET)
        add_text(slide, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), card_h - Inches(0.9),
                 text, size=17, color=WHITE, line_spacing=1.25)

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def demo_slide(prs, idx):
    slide = base_slide(prs)
    add_header(slide, "09 • LIVE EXAMPLE", "Real Example — Quote Review Prompt")

    x, y = Inches(0.7), Inches(2.55)
    w, h = Inches(11.933), Inches(4.0)
    add_rect(slide, x, y, w, h, fill=NAVY_900, line=CYAN, line_w=Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # terminal-style dots
    dot_y = y + Inches(0.3)
    for i, col in enumerate([RGBColor(0xFF, 0x5F, 0x57), RGBColor(0xFE, 0xBC, 0x2E), RGBColor(0x28, 0xC8, 0x40)]):
        add_rect(slide, x + Inches(0.35) + Inches(0.3) * i, dot_y, Inches(0.18), Inches(0.18),
                 fill=col, shape=MSO_SHAPE.OVAL)

    code = (
        "Role:   Senior Purchasing Analyst\n"
        "Task:   Review attached supplier quote\n\n"
        "Output:\n"
        "  - Pricing changes\n"
        "  - Lead time risks\n"
        "  - Missing information\n"
        "  - Recommended next action"
    )
    tb = slide.shapes.add_textbox(x + Inches(0.6), y + Inches(0.8),
                                  w - Inches(1.2), h - Inches(1.2))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        run = p.add_run()
        run.text = line
        run.font.name = MONO
        run.font.size = Pt(18)
        run.font.color.rgb = CYAN_L if line.startswith(("Role", "Task", "Output")) else INK_200
        run.font.bold = line.startswith(("Role", "Task", "Output"))

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def split_slide(prs, idx):
    slide = base_slide(prs)
    add_header(slide, "10 • BUILD ONCE, REUSE", "Building GPT Agents and Reusable Skills")

    y = Inches(2.6)
    w = Inches(5.85)
    h = Inches(2.6)
    for i, (title, desc) in enumerate([
        ("AGENTS", "Task-specific helpers tuned for a single recurring job."),
        ("SKILLS", "Reusable workflow playbooks applied across many tasks."),
    ]):
        x = Inches(0.7) + (w + Inches(0.2)) * i
        add_rect(slide, x, y, w, h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(1),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, x + Inches(0.4), y + Inches(0.35), w - Inches(0.8), Inches(0.4),
                 title, size=13, bold=True, color=CYAN_L, tracking=400)
        add_rect(slide, x + Inches(0.4), y + Inches(0.9), Inches(0.8), Emu(25400), fill=VIOLET)
        add_text(slide, x + Inches(0.4), y + Inches(1.1), w - Inches(0.8), h - Inches(1.3),
                 desc, size=18, color=WHITE, line_spacing=1.35)

    # footer strip
    strip_y = y + h + Inches(0.3)
    add_rect(slide, Inches(0.7), strip_y, Inches(11.933), Inches(0.75),
             fill=NAVY_700, line=CYAN, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, Inches(0.7), strip_y, Inches(11.933), Inches(0.75),
             "Quote review   •   PDF extraction   •   Quality summaries",
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             tracking=200)

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def divider_slide(prs, idx):
    slide = base_slide(prs)
    # stronger violet/cyan wash
    add_rect(slide, Inches(-1), Inches(-1), Inches(8), Inches(10), fill=VIOLET,
             shape=MSO_SHAPE.OVAL).fill.transparency = 0.82
    add_rect(slide, Inches(7), Inches(3), Inches(8), Inches(7), fill=CYAN,
             shape=MSO_SHAPE.OVAL).fill.transparency = 0.88

    add_portrait_placeholder(slide, Inches(0.7), Inches(1.0), Inches(4.2), Inches(5.5),
                             "SIDE PROFILE / FULL BODY")

    add_kicker(slide, Inches(5.4), Inches(2.5), "PART TWO")
    add_text(slide, Inches(5.4), Inches(3.05), Inches(7.5), Inches(1.2),
             "Operationalizing AI", size=48, bold=True, color=WHITE, line_spacing=1.05)
    add_text(slide, Inches(5.4), Inches(3.95), Inches(7.5), Inches(1.2),
             "at Accutite", size=48, bold=True, color=VIOLET_L, line_spacing=1.05)

    add_rect(slide, Inches(5.4), Inches(5.05), Inches(1.3), Emu(50800), fill=CYAN)
    add_text(slide, Inches(5.4), Inches(5.2), Inches(7), Inches(0.5),
             "FROM AWARENESS TO EXECUTION", size=11, bold=True, color=INK_300, tracking=400)

    # watermark logo
    add_text(slide, Inches(8.5), Inches(6.7), Inches(4), Inches(0.4),
             "ACCUTITE", size=14, bold=True, color=INK_400, tracking=500,
             align=PP_ALIGN.RIGHT)
    add_page_num(slide, idx)


def steps_slide(prs, idx):
    slide = base_slide(prs)
    add_header(slide, "11 • ROADMAP", "Automation Roadmap: Start Small, Win Fast")

    steps = [
        ("01", "Standardize prompts"),
        ("02", "Build reusable prompt libraries and skills"),
        ("03", "Automate report analysis"),
        ("04", "Expand into agent workflows"),
    ]
    w = Inches(2.88)
    h = Inches(2.9)
    y = Inches(2.6)
    for i, (n, text) in enumerate(steps):
        x = Inches(0.7) + (w + Inches(0.15)) * i
        add_rect(slide, x, y, w, h, fill=NAVY_700, line=VIOLET_L, line_w=Pt(1),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, x, y + Inches(0.4), w, Inches(0.4),
                 n, size=12, bold=True, color=CYAN_L, tracking=400, align=PP_ALIGN.CENTER)
        add_rect(slide, x + w / 2 - Inches(0.35), y + Inches(0.95), Inches(0.7), Emu(25400),
                 fill=VIOLET)
        add_text(slide, x + Inches(0.3), y + Inches(1.2), w - Inches(0.6), h - Inches(1.4),
                 text, size=16, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        # arrow between steps
        if i < len(steps) - 1:
            ax = x + w + Inches(0.01)
            add_text(slide, ax, y + h / 2 - Inches(0.25), Inches(0.15), Inches(0.5),
                     "›", size=20, bold=True, color=VIOLET_L,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def beforeafter_slide(prs, idx):
    slide = base_slide(prs)
    add_header(slide, "13 • IMPACT", "From Reports to Decisions")

    cols = [
        ("BEFORE", "Manual review", INK_300, VIOLET_L),
        ("AFTER", "AI summary and flags", CYAN_L, CYAN),
        ("OUTCOME", "Faster decisions,\nfewer misses", VIOLET_L, VIOLET),
    ]
    w = Inches(3.75)
    gap = Inches(0.15)
    h = Inches(3.2)
    y = Inches(2.8)
    for i, (label, body, label_color, border) in enumerate(cols):
        x = Inches(0.7) + (w + gap + Inches(0.25)) * i
        add_rect(slide, x, y, w, h, fill=NAVY_700, line=border, line_w=Pt(1.5),
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, x, y + Inches(0.45), w, Inches(0.4),
                 label, size=12, bold=True, color=label_color, tracking=400,
                 align=PP_ALIGN.CENTER)
        add_rect(slide, x + w / 2 - Inches(0.4), y + Inches(1.0), Inches(0.8), Emu(25400), fill=CYAN)
        add_text(slide, x + Inches(0.4), y + Inches(1.2), w - Inches(0.8), h - Inches(1.3),
                 body, size=20, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3, bold=True)
        if i < len(cols) - 1:
            ax = x + w + Inches(0.05)
            add_text(slide, ax, y + h / 2 - Inches(0.3), Inches(0.3), Inches(0.6),
                     "→", size=28, bold=True, color=VIOLET_L,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_page_num(slide, idx)
    add_footer(slide)
    add_logo_lockup(slide, Inches(11.2), Inches(0.35))


def closing_slide(prs, idx):
    slide = base_slide(prs)

    add_portrait_placeholder(slide, Inches(0.7), Inches(1.0), Inches(4.2), Inches(5.5),
                             "CLOSING PORTRAIT")

    add_kicker(slide, Inches(5.4), Inches(1.15), "CALL TO ACTION")
    add_text(slide, Inches(5.4), Inches(1.7), Inches(7.5), Inches(1.3),
             "Become Operators,", size=48, bold=True, color=WHITE, line_spacing=1.0)
    add_text(slide, Inches(5.4), Inches(2.6), Inches(7.5), Inches(1.3),
             "Not Observers.", size=48, bold=True, color=CYAN_L, line_spacing=1.0)
    add_rect(slide, Inches(5.4), Inches(3.7), Inches(1.2), Emu(50800), fill=VIOLET)

    # signature card
    sx, sy = Inches(5.4), Inches(4.0)
    sw, sh = Inches(7.2), Inches(2.2)
    add_rect(slide, sx, sy, sw, sh, fill=NAVY_700, line=VIOLET_L, line_w=Pt(1),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    rows = [
        ("LED BY", "Levee Hedrick"),
        ("ROLE", "Material Control Manager"),
        ("MISSION", "Driving AI adoption across Purchasing & Quality"),
    ]
    for i, (label, val) in enumerate(rows):
        ry = sy + Inches(0.25) + Inches(0.6) * i
        add_text(slide, sx + Inches(0.4), ry, Inches(1.6), Inches(0.4),
                 label, size=10, bold=True, color=CYAN_L, tracking=400,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, sx + Inches(2.1), ry, sw - Inches(2.5), Inches(0.4),
                 val, size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_logo_lockup(slide, Inches(10.7), Inches(6.6), large=True)
    add_page_num(slide, idx)


# ---------- main ----------
def build(output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1
    slide_title(prs)
    # 2
    content_slide(prs, 2, "01 • VISION", "The New Operating System for Fastener Work", [
        "ChatGPT is a daily work amplifier for the fastener industry.",
        "Focus on speed, consistency, and better decisions.",
        "Shift from casual use to disciplined adoption.",
    ])
    # 3
    content_slide(prs, 3, "02 • CONTEXT", "Why This Matters at Accutite Right Now", [
        "AI understanding is inconsistent across the team.",
        "Standardization matters more than random LLM experimentation.",
        "Paid accounts should create measurable operational value.",
    ])
    # 4
    chips_slide(prs, 4)
    # 5
    cards_slide(prs, 5, "04 • DEPARTMENT FOCUS", "Purchasing Department Use Cases", [
        "Supplier email review",
        "Quote comparison",
        "PDF data extraction",
        "Supplier follow-up drafting",
    ])
    # 6
    cards_slide(prs, 6, "05 • DEPARTMENT FOCUS", "Quality Department Use Cases", [
        "Summarize inspection and corrective action documents",
        "Review drawings and blueprint notes",
        "Compare document versions",
        "Generate action checklists",
    ])
    # 7
    content_slide(prs, 7, "06 • TECHNICAL DOCUMENTS", "Working with Drawings, Blueprints, and PDFs", [
        "Summarize technical documents at a glance.",
        "Extract revision notes and requirements.",
        "Support triage before expert approval.",
        "Preserve human technical authority.",
    ])
    # 8
    content_slide(prs, 8, "07 • ANALYTICS", "From Reports to KPIs Without ERP Integration", [
        "Analyze exported reports directly in ChatGPT.",
        "Summarize KPI trends across periods.",
        "Detect exceptions and outliers.",
        "Create management-ready insights.",
    ])
    # 9
    content_slide(prs, 9, "08 • DISCIPLINE", "Prompt Engineering for the Fastener Industry", [
        "Define role, task, context, format, and criteria.",
        "Use repeatable prompting frameworks.",
        "Emphasize quality and consistency over cleverness.",
    ])
    # 10
    demo_slide(prs, 10)
    # 11
    split_slide(prs, 11)
    # 12
    divider_slide(prs, 12)
    # 13
    steps_slide(prs, 13)
    # 14
    content_slide(prs, 14, "12 • 30-DAY TARGET", "What Good Looks Like in 30 Days", [
        "One defined workflow per user.",
        "Shared prompt wins across the team.",
        "Repeatable processes in Purchasing and Quality.",
        "Better leadership reporting.",
    ])
    # 15
    beforeafter_slide(prs, 15)
    # 16
    closing_slide(prs, 16)

    prs.save(output)
    return output


if __name__ == "__main__":
    out = Path(__file__).parent / "Accutite_ChatGPT_Fastener_Deck.pptx"
    path = build(out)
    print(f"Wrote {path} ({path.stat().st_size / 1024:.1f} KB)")
